from matplotlib.widgets import Slider, Button
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import av
import matplotlib.pyplot as plt
import matplotlib.animation as animation
import numpy as np
import os
import pandas as pd
import time

# Start time
t_start = time.time()

# Create a list of DataFrame subsets. Each subset contains av.con rows from Df_dataset
# The subsets are created by slicing Df_dataset at intervals of av.con rows
av.Df_dataset_subsets = [av.Df_dataset.iloc[i*av.con:(i+1)*av.con] for i in range(av.con)]
# Prepare the points for each subset
points_x = [[av.Df_dataset_subset[av.Df_dataset_subset['poiID'] == poi_id]['x'].tolist() for poi_id in range(av.poi)] for av.Df_dataset_subset in av.Df_dataset_subsets]
points_y = [[av.Df_dataset_subset[av.Df_dataset_subset['poiID'] == poi_id]['y'].tolist() for poi_id in range(av.poi)] for av.Df_dataset_subset in av.Df_dataset_subsets]
points_timestamps = [[av.Df_dataset_subset[av.Df_dataset_subset['poiID'] == poi_id]['tstID'].tolist() for poi_id in range(av.poi)] for av.Df_dataset_subset in av.Df_dataset_subsets]

points_x_interpolated = [
    [
        np.interp(
            np.linspace(min(timestamps), max(timestamps), av.num_frames),
            timestamps,
            point_x,
        )
        for timestamps, point_x in zip(points_timestamps[i], points_x[i])
    ]
    for i in range(av.con)
]

points_y_interpolated = [
    [
        np.interp(
            np.linspace(min(timestamps), max(timestamps), av.num_frames),
            timestamps,
            point_y,
        )
        for timestamps, point_y in zip(points_timestamps[i], points_y[i])
    ]
    for i in range(av.con)
]

rows = av.con // 3 + 1

# Create the figure and axes
fig, axs = plt.subplots(rows, 3, sharex='all', sharey='all', figsize=(12, 7))  # 2 rows, 3 columns

# Dynamically create the titles
titles = ["Title for Subfigure {}".format(i+1) for i in range(av.con)]

# Iterate over the original 2D array of axes
counter = 0
for i in range(2):
    for j in range(3):
        if counter < len(titles):
            ax = axs[i, j]
            ax.set_title(titles[i*3 + j])
            ax.set_xlim(av.min_boundary_x, av.max_boundary_x)  # Adjust the x-axis limits based on the dataset
            ax.set_ylim(av.min_boundary_y, av.max_boundary_y)  # Adjust the y-axis limits based on the dataset
            counter = counter + 1

colors = ['black', 'red', 'red']  # Add more colors if needed

points = []
for i in range(av.con):
    # print(i, av.poi, i // 3, i % 3)
    ax = axs[i // 3, i % 3]  # Get the corresponding axis from the original 2D array
    point_set = []
    for j in range(av.poi):
        point = ax.plot([], [], "o", color=colors[j % len(colors)])[0]
        point_set.append(point)
    points.append(point_set)

def init():
    for point_set in points:
        for point in point_set:
            point.set_data([], [])
    return [point for point_set in points for point in point_set]

def update(frame):
    for i, point_set in enumerate(points):
        for j, point in enumerate(point_set):
            x = points_x_interpolated[i][j][frame]
            y = points_y_interpolated[i][j][frame]
            point.set_data([x], [y])  # Convert x and y to sequences
    return [point for point_set in points for point in point_set]

# Declare the animation object globally
ani = None

def setup_animation():
    # Create the animation without starting it immediately
    global ani
    ani = animation.FuncAnimation(
        fig,
        update,
        frames=av.num_frames,
        init_func=init,
        blit=True,
        interval=50,
        repeat=False,
    )
    return ani

# Create a slider to control the animation
slider_ax = plt.axes([0.15, 0.02, 0.7, 0.03])
slider = Slider(slider_ax, "Frame", 0, av.num_frames - 1, valinit=0, valstep=1)

def slider_update(val):
    frame = int(val)
    update(frame)
    plt.draw()

slider.on_changed(slider_update)

# Create a button to play the animation again
def restart_animation(event):
    global ani
    if ani is not None:
        if ani.event_source is not None:  # Check that event_source is not None
            ani.event_source.stop()  # Stop the previous animation
    ani = setup_animation()
    ani.event_source.start()  # Start the new animation
    
button_ax = plt.axes([0.8, 0.025, 0.1, 0.04])
button = Button(button_ax, 'Restart')
button.on_clicked(restart_animation)

# Setup the initial animation
ani = setup_animation()

# Save the animation as an mp4 file
video_name = "animation.mp4"
ani.save(video_name, writer='ffmpeg' , dpi = 50)

# Save the animation as an HTML file
html_name = "animation.html"
ani.save(html_name, writer='html', dpi = 100)

# Create a PowerPoint presentation
presentation = Presentation()

# Create a slide with the animated visualization
slide_layout = presentation.slide_layouts[6]  # Choose a layout for the slide
slide = presentation.slides.add_slide(slide_layout)

# Set the size and position of the video frame on the slide
left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(6)

# Add a title to the slide
title_text = "Animated Visualization"  # Set the title text
title_shape = slide.shapes.add_textbox(left, Inches(0.1), width, Inches(0.5))  # Adjust the position and height of the title
title_frame = title_shape.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title_text
title_paragraph.alignment = PP_ALIGN.CENTER  # Adjust the alignment of the title
title_paragraph.font.name = "Arial"  # Change the font to Arial
title_paragraph.font.size = Pt(32)  # Adjust the font size of the title

# Adjust the paragraph spacing
title_paragraph.space_before = Pt(0)
title_paragraph.space_after = Pt(0)

# Add the video to the slide
video_path = "animation.mp4"
slide.shapes.add_movie(video_path, left, top, width, height)

# Save the PowerPoint presentation
pptx_name = "dynamics.pptx"
presentation.save(pptx_name)

plt.clf()  # Clear the figure to start with a new blank figure

# End and print time
print('Time elapsed for running module "N_VA_DynamicAbsolute": {:.3f} sec.'.format(time.time() - t_start))